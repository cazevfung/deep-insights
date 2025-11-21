"""Article content scraper with Playwright and trafilatura."""
import time
import uuid
from datetime import datetime
from typing import Dict
from loguru import logger
from playwright.sync_api import Page
import trafilatura
from urllib.parse import unquote
from scrapers.base_scraper import BaseScraper


class ArticleScraper(BaseScraper):
    """Extract text content from articles and web pages."""
    
    def __init__(self, **kwargs):
        """Initialize article scraper."""
        super().__init__(**kwargs)
        self.method_preference = self.scraper_config.get('method_preference', 'playwright')
        self.min_content_words = self.scraper_config.get('min_content_words', 50)
        self.remove_blocking = self.scraper_config.get('remove_blocking_elements', True)
    
    def validate_url(self, url: str) -> bool:
        """
        Check if URL is valid (accepts any http/https URL or file:// URL except Reddit).
        
        Args:
            url: URL to validate
            
        Returns:
            True if valid URL
        """
        if url.startswith('file://'):
            return True
        return url.startswith(('http://', 'https://')) and 'reddit.com' not in url
    
    def extract(self, url: str, batch_id: str = None, link_id: str = None) -> Dict:
        """
        Extract content from article.
        Tries Playwright first, falls back to trafilatura.
        Supports file:// URLs for local files (backward compatible).
        
        Args:
            url: Article URL or file:// URL
            
        Returns:
            Dictionary with extraction results
        """
        start_time = time.time()
        result = None
        
        # Set extraction context for progress reporting
        self._set_extraction_context(batch_id, link_id, url)
        
        # Handle file:// URLs for local files (NEW feature, backward compatible)
        # Graceful fallback: If file extraction fails, continue with normal flow
        if url.startswith('file://'):
            try:
                result = self._extract_from_local_file(url)
                if result.get('success'):
                    elapsed_time = round(time.time() - start_time, 2)
                    logger.info(f"[Article] Local file extraction completed in {elapsed_time}s - Success: {result['success']}")
                    return result
                else:
                    # File extraction failed, log and continue with normal flow
                    logger.warning(f"[Article] Local file extraction failed: {result.get('error', 'Unknown error')}, falling back to normal extraction")
                    result = None  # Reset to allow normal flow
            except Exception as e:
                # Graceful fallback: If file extraction throws exception, continue with normal flow
                logger.warning(f"[Article] Exception during local file extraction: {e}, falling back to normal extraction")
                result = None  # Reset to allow normal flow
        
        # Try Playwright first if preference is playwright
        if self.method_preference == 'playwright':
            result = self._extract_with_playwright(url)
            
            # If Playwright fails and we got content but it's too short, try trafilatura
            if not result['success'] or (result.get('content') and len(result.get('content', '').split()) < self.min_content_words):
                logger.info("[Article] Falling back to trafilatura...")
                result = self._extract_with_trafilatura(url)
        
        # Try trafilatura if preference is trafilatura or as fallback
        if self.method_preference == 'trafilatura' or result is None:
            result = self._extract_with_trafilatura(url)
        
        # Ensure all required fields are present
        if 'extraction_timestamp' not in result:
            result['extraction_timestamp'] = datetime.now().isoformat()
        
        # Assign a random identifier to this article extraction if not present
        try:
            if 'article_id' not in result or not result.get('article_id'):
                result['article_id'] = uuid.uuid4().hex[:12]
        except Exception:
            # Best-effort; do not fail extraction due to ID assignment
            result['article_id'] = uuid.uuid4().hex[:12]

        elapsed_time = round(time.time() - start_time, 2)
        logger.info(f"[Article] Extraction completed in {elapsed_time}s - Success: {result['success']}")
        
        return result
    
    def _extract_from_local_file(self, file_url: str) -> Dict:
        """
        Extract content from a local file using file:// URL.
        
        This is a NEW feature added for backward compatibility.
        Returns result dict matching existing extract() return format.
        
        Args:
            file_url: file:// URL pointing to local file
            
        Returns:
            Dictionary with extraction results matching existing format:
            {
                'success': bool,
                'content': str,
                'metadata': dict,
                'extraction_timestamp': str,
                'article_id': str,
                'error': str (if failed)
            }
        """
        from pathlib import Path
        
        try:
            def _build_success_result(content: str, extraction_method: str, extra_metadata: Dict = None) -> Dict:
                """
                Build a standardized success payload for local file extractions.
                Ensures batch/link identifiers are preserved for downstream routing.
                """
                if not content or not content.strip():
                    raise ValueError(f"No textual content could be extracted from {path.name}")
                
                cleaned_content = content.strip()
                word_count = len(cleaned_content.split())
                file_size = path.stat().st_size if path.exists() else 0
                
                metadata = {
                    'title': path.stem,
                    'source': 'local_file',
                    'url': file_url,
                    'file_path': str(path),
                    'file_size': file_size,
                    'word_count': word_count,
                    'file_extension': path.suffix,
                    'extraction_method': extraction_method
                }
                if extra_metadata:
                    metadata.update(extra_metadata)
                
                return {
                    'success': True,
                    'url': file_url,
                    'content': cleaned_content,
                    'title': metadata.get('title', ''),
                    'author': metadata.get('author', ''),
                    'publish_date': metadata.get('publish_date', ''),
                    'source': metadata.get('source', 'local_file'),
                    'language': metadata.get('language', 'auto'),
                    'word_count': word_count,
                    'extraction_method': extraction_method,
                    'extraction_timestamp': datetime.now().isoformat(),
                    'batch_id': self._current_batch_id,
                    'link_id': self._current_link_id,
                    'metadata': metadata,
                    'article_id': uuid.uuid4().hex[:12],
                    'error': None
                }
        
            # Parse file:// URL
            # file:///C:/path/to/file.md or file:///path/to/file.md
            file_path = file_url.replace('file:///', '').replace('file://', '')
            file_path = unquote(file_path)
            
            # Handle Windows paths (file:///C:/...)
            if file_path.startswith('/') and len(file_path) > 1 and file_path[1] == ':':
                file_path = file_path[1:]  # Remove leading /
            
            path = Path(file_path)
            
            if not path.exists():
                logger.error(f"[Article] Local file not found: {file_path}")
                return {
                    'success': False,
                    'error': f'File not found: {file_path}',
                    'content': '',
                    'metadata': {
                        'title': path.stem,
                        'source': 'local_file',
                        'url': file_url,
                        'file_path': str(path)
                    },
                    'extraction_timestamp': datetime.now().isoformat(),
                    'article_id': uuid.uuid4().hex[:12]
                }
            
            # Handle different file types
            suffix = path.suffix.lower()
            
            # PDF files - convert to markdown using PDF converter
            if suffix == '.pdf':
                self._report_progress("loading", 30, f"Converting PDF: {path.name}")
                logger.info(f"[Article] Processing PDF file: {path}")
                
                # Suppress verbose DEBUG logging from PDF parsing libraries
                import logging
                pdf_loggers = ['pdfminer', 'pdfminer.pdfparser', 'pdfminer.pdfdocument', 
                              'pdfminer.pdfinterp', 'pdfminer.pdfpage', 'pdfminer.converter',
                              'pdfminer.layout', 'pdfminer.psparser', 'pdfminer.cmapdb',
                              'pdfminer.six', 'pypdf', 'PyPDF2']
                
                # Store original log levels
                original_levels = {}
                root_logger = logging.getLogger()
                original_root_level = root_logger.level
                
                # Suppress PDF library logs
                root_logger.setLevel(logging.WARNING)
                for logger_name in pdf_loggers:
                    pdf_logger = logging.getLogger(logger_name)
                    original_levels[logger_name] = pdf_logger.level
                    pdf_logger.setLevel(logging.WARNING)
                
                try:
                    from utils.pdf_to_markdown import convert_pdf_to_markdown
                    content = convert_pdf_to_markdown(path)
                    extraction_method = 'pdf_to_markdown'
                except ImportError:
                    logger.warning("[Article] PDF to markdown converter not available, using pypdf fallback")
                    from pypdf import PdfReader
                    reader = PdfReader(str(path))
                    pages = [page.extract_text() or "" for page in reader.pages]
                    content = "\n\n".join(pages)
                    extraction_method = 'pypdf_fallback'
                except Exception as e:
                    logger.error(f"[Article] PDF conversion failed: {e}", exc_info=True)
                    # Try pypdf fallback
                    try:
                        from pypdf import PdfReader
                        reader = PdfReader(str(path))
                        pages = [page.extract_text() or "" for page in reader.pages]
                        content = "\n\n".join(pages)
                        extraction_method = 'pypdf_fallback'
                    except Exception as fallback_error:
                        logger.error(f"[Article] PDF extraction (fallback) also failed: {fallback_error}", exc_info=True)
                        raise ValueError(f"Failed to extract text from PDF: {str(e)}") from e
                finally:
                    # Restore original log levels
                    root_logger.setLevel(original_root_level)
                    for logger_name, level in original_levels.items():
                        logging.getLogger(logger_name).setLevel(level)
                
                self._report_progress("extracting", 90, f"Extracted {len(content.split())} words from PDF {path.name}")
                
                return _build_success_result(
                    content,
                    extraction_method,
                    extra_metadata={'pdf_total_pages': None}
                )
            
            # Word documents (.doc, .docx)
            elif suffix in ['.doc', '.docx']:
                self._report_progress("loading", 30, f"Extracting from Word document: {path.name}")
                logger.info(f"[Article] Processing Word document: {path}")
                
                try:
                    from docx import Document
                    doc = Document(str(path))
                    content = "\n".join(p.text for p in doc.paragraphs)
                    extraction_method = 'python_docx'
                except Exception as e:
                    logger.error(f"[Article] Word document extraction failed: {e}", exc_info=True)
                    raise ValueError(f"Failed to extract text from Word document: {str(e)}") from e
                
                self._report_progress("extracting", 90, f"Extracted {len(content.split())} words from Word document {path.name}")
                
                return _build_success_result(content, extraction_method)
            
            # PowerPoint files (.ppt, .pptx)
            elif suffix in ['.ppt', '.pptx']:
                self._report_progress("loading", 30, f"Extracting from PowerPoint: {path.name}")
                logger.info(f"[Article] Processing PowerPoint file: {path}")
                
                try:
                    from pptx import Presentation
                    prs = Presentation(str(path))
                    texts = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                texts.append(shape.text)
                    content = "\n".join(texts)
                    extraction_method = 'python_pptx'
                except Exception as e:
                    logger.error(f"[Article] PowerPoint extraction failed: {e}", exc_info=True)
                    raise ValueError(f"Failed to extract text from PowerPoint: {str(e)}") from e
                
                self._report_progress("extracting", 90, f"Extracted {len(content.split())} words from PowerPoint {path.name}")
                
                return _build_success_result(content, extraction_method)
            
            # Excel files (.xls, .xlsx) - convert to JSON
            elif suffix in ['.xls', '.xlsx']:
                self._report_progress("loading", 30, f"Extracting from Excel: {path.name}")
                logger.info(f"[Article] Processing Excel file: {path}")
                
                try:
                    import pandas as pd
                    sheets = pd.read_excel(str(path), sheet_name=None, engine="openpyxl")
                    content_parts = []
                    for sheet_name, df in sheets.items():
                        sanitized = df.dropna(how="all").dropna(axis=1, how="all").fillna("")
                        content_parts.append(f"Sheet: {sheet_name}\n")
                        content_parts.append(sanitized.to_string())
                        content_parts.append("\n\n")
                    content = "\n".join(content_parts)
                    extraction_method = 'pandas_excel'
                except Exception as e:
                    logger.error(f"[Article] Excel extraction failed: {e}", exc_info=True)
                    raise ValueError(f"Failed to extract text from Excel: {str(e)}") from e
                
                self._report_progress("extracting", 90, f"Extracted {len(content.split())} words from Excel {path.name}")
                
                return _build_success_result(content, extraction_method)
            
            # Plain text files (.txt, .md, .markdown)
            else:
                self._report_progress("loading", 50, f"Reading file: {path.name}")
                
                # Try UTF-8 first, fallback to other encodings
                try:
                    content = path.read_text(encoding='utf-8')
                except UnicodeDecodeError:
                    # Graceful fallback: Try common encodings
                    for encoding in ['utf-8-sig', 'latin-1', 'cp1252']:
                        try:
                            content = path.read_text(encoding=encoding)
                            logger.info(f"[Article] Successfully read file with {encoding} encoding")
                            break
                        except UnicodeDecodeError:
                            continue
                    else:
                        # All encodings failed
                        raise UnicodeDecodeError('utf-8', b'', 0, 1, 'Could not decode file with any encoding')
                
                # Extract metadata (matching existing format)
                self._report_progress("extracting", 90, f"Extracted {len(content.split())} words from {path.name}")
                
                return _build_success_result(content, 'local_file_read')
            
        except UnicodeDecodeError as e:
            logger.error(f"[Article] Failed to decode file {file_url}: {e}")
            return {
                'success': False,
                'error': f'Failed to decode file: {str(e)}',
                'content': '',
                'metadata': {
                    'title': path.stem if 'path' in locals() else 'unknown',
                    'source': 'local_file',
                    'url': file_url
                },
                'extraction_timestamp': datetime.now().isoformat(),
                'article_id': uuid.uuid4().hex[:12],
                'batch_id': self._current_batch_id,
                'link_id': self._current_link_id
            }
        except Exception as e:
            logger.error(f"[Article] Failed to read local file {file_url}: {e}", exc_info=True)
            return {
                'success': False,
                'error': f'Failed to read file: {str(e)}',
                'content': '',
                'metadata': {
                    'title': path.stem if 'path' in locals() else 'unknown',
                    'source': 'local_file',
                    'url': file_url
                },
                'extraction_timestamp': datetime.now().isoformat(),
                'article_id': uuid.uuid4().hex[:12],
                'batch_id': self._current_batch_id,
                'link_id': self._current_link_id
            }
    
    def _extract_with_playwright(self, url: str) -> Dict:
        """
        Extract content using Playwright.
        
        Args:
            url: Article URL
        
        Returns:
            Dictionary with extraction results
        """
        try:
            context = self._create_context()
            page = context.new_page()
            
            # Navigate
            self._report_progress("loading", 10, "Loading article")
            page.goto(url, wait_until='networkidle', timeout=self.timeout)
            
            # Check for cancellation
            if self._check_cancelled():
                logger.info(f"[Article] Cancellation detected, force closing browser for {url}")
                try:
                    page.close()
                except:
                    pass
                try:
                    context.close()
                except:
                    pass
                self.close(force_kill=True)  # Force kill browser processes
                return self._error_result(url, "Cancelled by user", batch_id, link_id)
            
            time.sleep(2.0)
            self._report_progress("loading", 30, "Article loaded")
            
            # Check for cancellation again
            if self._check_cancelled():
                logger.info(f"[Article] Cancellation detected, force closing browser for {url}")
                try:
                    page.close()
                except:
                    pass
                try:
                    context.close()
                except:
                    pass
                self.close(force_kill=True)  # Force kill browser processes
                return self._error_result(url, "Cancelled by user", batch_id, link_id)
            
            # Scroll to trigger lazy loading
            page.evaluate("window.scrollTo(0, document.body.scrollHeight)")
            time.sleep(0.5)
            self._report_progress("loading", 40, "Loading additional content")
            
            # Click expand buttons
            self._click_expand_buttons(page)
            self._report_progress("extracting", 50, "Expanding content")
            
            # Scroll again
            page.evaluate("window.scrollTo(0, document.body.scrollHeight)")
            time.sleep(0.5)
            
            # Remove blocking elements if enabled
            if self.remove_blocking:
                self._remove_blocking_elements(page)
            
            # Extract metadata
            self._report_progress("extracting", 60, "Extracting metadata")
            metadata = self._extract_metadata(page, url)
            
            # Extract content using multiple selectors
            self._report_progress("extracting", 70, "Extracting article content")
            content_selectors = [
                'article',
                'main',
                '[role="article"]',
                '.article-content',
                '.content-wrapper',
                '.post-content',
                '.entry-content',
                '#content',
                '.article-body',
                '.news-content',
                '.article-text',
                'div[class*="content"]',
                'div[class*="article"]',
            ]
            
            content = None
            for selector in content_selectors:
                try:
                    element = page.locator(selector).first
                    if element.count() > 0:
                        content = element.inner_text(timeout=2000)
                        if content and len(content.split()) > self.min_content_words:
                            logger.debug(f"[Article] Found content with selector: {selector}")
                            break
                except:
                    continue
            
            # Fallback to body if no specific content found
            if not content or len(content.split()) < self.min_content_words:
                try:
                    page.evaluate("""
                        document.querySelectorAll('nav, header, footer, aside, .ad, .advertisement, .comment').forEach(el => el.remove());
                    """)
                    body = page.locator('body')
                    content = body.inner_text(timeout=2000)
                except:
                    pass
            
            # Clean up
            page.close()
            context.close()
            
            if content:
                content = self._clean_content(content)
                word_count = len(content.split())
                self._report_progress("extracting", 100, f"Extracted {word_count} words")
                
                if word_count >= self.min_content_words:
                    return {
                        'success': True,
                        'url': url,
                        'content': content,
                        'title': metadata.get('title', ''),
                        'author': metadata.get('author', ''),
                        'publish_date': metadata.get('publish_date', ''),
                        'source': metadata.get('source', ''),
                        'language': 'auto',
                        'word_count': word_count,
                        'extraction_method': 'article_playwright',
                        'extraction_timestamp': datetime.now().isoformat(),
                        'batch_id': self._current_batch_id,
                        'link_id': self._current_link_id,
                        'error': None
                    }
                else:
                    return {
                        'success': False,
                        'url': url,
                        'content': content,
                        'title': metadata.get('title', ''),
                        'author': metadata.get('author', ''),
                        'publish_date': metadata.get('publish_date', ''),
                        'source': metadata.get('source', ''),
                        'language': 'auto',
                        'word_count': word_count,
                        'extraction_method': 'article_playwright',
                        'extraction_timestamp': datetime.now().isoformat(),
                        'batch_id': self._current_batch_id,
                        'link_id': self._current_link_id,
                        'error': f'Content too short ({word_count} words)'
                    }
            else:
                return {
                    'success': False,
                    'url': url,
                    'content': None,
                    'title': metadata.get('title', ''),
                    'author': metadata.get('author', ''),
                    'publish_date': metadata.get('publish_date', ''),
                    'source': metadata.get('source', ''),
                    'language': 'auto',
                    'word_count': 0,
                    'extraction_method': 'article_playwright',
                    'extraction_timestamp': datetime.now().isoformat(),
                    'batch_id': self._current_batch_id,
                    'link_id': self._current_link_id,
                    'error': 'No content found'
                }
        
        except Exception as e:
            logger.error(f"[Article] Playwright extraction failed: {e}")
            return {
                'success': False,
                'url': url,
                'content': None,
                'title': '',
                'author': '',
                'publish_date': '',
                'source': '',
                'language': 'auto',
                'word_count': 0,
                'extraction_method': 'article_playwright',
                'extraction_timestamp': datetime.now().isoformat(),
                'batch_id': self._current_batch_id,
                'link_id': self._current_link_id,
                'error': str(e)
            }
    
    def _extract_with_trafilatura(self, url: str) -> Dict:
        """
        Extract content using trafilatura.
        
        Args:
            url: Article URL
        
        Returns:
            Dictionary with extraction results
        """
        try:
            logger.debug("[Article] Using trafilatura...")
            
            # Download and extract
            downloaded = trafilatura.fetch_url(url)
            if not downloaded:
                return {
                    'success': False,
                    'url': url,
                    'content': None,
                    'title': '',
                    'author': '',
                    'publish_date': '',
                    'source': self._extract_domain(url),
                    'language': 'auto',
                    'word_count': 0,
                    'extraction_method': 'article_trafilatura',
                    'extraction_timestamp': datetime.now().isoformat(),
                    'error': 'Download failed'
                }
            
            # Extract content
            try:
                content = trafilatura.extract(
                    downloaded,
                    include_comments=False,
                    include_tables=True,
                    include_links=False,
                    output_format='txt'
                )
            except Exception as e:
                logger.debug(f"Trafilatura extract failed: {e}")
                content = None
            
            if content:
                content = self._clean_content(content)
                word_count = len(content.split())
                
                if word_count >= self.min_content_words:
                    # Extract metadata from downloaded HTML
                    try:
                        metadata = trafilatura.extract_metadata(downloaded)
                        if metadata and hasattr(metadata, 'get'):
                            metadata_dict = metadata
                        else:
                            metadata_dict = None
                    except:
                        metadata_dict = None
                    
                    return {
                        'success': True,
                        'url': url,
                        'content': content,
                        'title': metadata_dict.get('title', '') if metadata_dict and hasattr(metadata_dict, 'get') else '',
                        'author': metadata_dict.get('author', '') if metadata_dict and hasattr(metadata_dict, 'get') else '',
                        'publish_date': metadata_dict.get('date', '') if metadata_dict and hasattr(metadata_dict, 'get') else '',
                        'source': self._extract_domain(url),
                        'language': 'auto',
                        'word_count': word_count,
                        'extraction_method': 'article_trafilatura',
                        'extraction_timestamp': datetime.now().isoformat(),
                        'batch_id': self._current_batch_id,
                        'link_id': self._current_link_id,
                        'error': None
                    }
                else:
                    return {
                        'success': False,
                        'url': url,
                        'content': content,
                        'title': '',
                        'author': '',
                        'publish_date': '',
                        'source': self._extract_domain(url),
                        'language': 'auto',
                        'word_count': word_count,
                        'extraction_method': 'article_trafilatura',
                        'extraction_timestamp': datetime.now().isoformat(),
                        'batch_id': self._current_batch_id,
                        'link_id': self._current_link_id,
                        'error': f'Content too short ({word_count} words)'
                    }
            else:
                return {
                    'success': False,
                    'url': url,
                    'content': None,
                    'title': '',
                    'author': '',
                    'publish_date': '',
                    'source': self._extract_domain(url),
                    'language': 'auto',
                    'word_count': 0,
                    'extraction_method': 'article_trafilatura',
                    'extraction_timestamp': datetime.now().isoformat(),
                    'batch_id': self._current_batch_id,
                    'link_id': self._current_link_id,
                    'error': 'Extraction failed'
                }
        
        except Exception as e:
            logger.error(f"[Article] Trafilatura extraction failed: {e}")
            return {
                'success': False,
                'url': url,
                'content': None,
                'title': '',
                'author': '',
                'publish_date': '',
                'source': self._extract_domain(url),
                'language': 'auto',
                'word_count': 0,
                'extraction_method': 'article_trafilatura',
                'extraction_timestamp': datetime.now().isoformat(),
                'batch_id': self._current_batch_id,
                'link_id': self._current_link_id,
                'error': str(e)
            }
    
    def _click_expand_buttons(self, page: Page):
        """Click 'Read More' / 'Expand' buttons."""
        expand_selectors = [
            'button:has-text("Read More")',
            'button:has-text("Read more")',
            'button:has-text("Show More")',
            'button:has-text("展开全文")',
            'button:has-text("展开")',
            'button:has-text("查看更多")',
            'a:has-text("Read More")',
            '.read-more',
            '.show-more'
        ]
        
        for selector in expand_selectors:
            try:
                buttons = page.locator(selector).all()
                for button in buttons[:3]:  # Limit to 3 buttons
                    try:
                        if button.is_visible(timeout=1000):
                            button.click()
                            time.sleep(0.5)
                    except:
                        continue
            except:
                continue
    
    def _remove_blocking_elements(self, page: Page):
        """Remove paywalls, overlays, modals."""
        try:
            page.evaluate("""
                const selectors = [
                    '.paywall', '.paywall-overlay', '.subscription-required',
                    '[class*="paywall"]', '[class*="overlay"]', '[class*="modal"]',
                    '.cookie-notice', '.gdpr-banner', '.privacy-notice',
                    '.newsletter', '.email-signup', '.subscribe-modal'
                ];
                
                selectors.forEach(selector => {
                    try {
                        document.querySelectorAll(selector).forEach(el => {
                            if (el.getBoundingClientRect().width > 100 && 
                                el.getBoundingClientRect().height > 100) {
                                el.remove();
                            }
                        });
                    } catch(e) {}
                });
            """)
        except:
            pass
    
    def _extract_metadata(self, page: Page, url: str) -> Dict:
        """Extract metadata from page."""
        metadata = {
            'title': '',
            'author': '',
            'publish_date': '',
            'source': self._extract_domain(url)
        }
        
        # Title
        title_selectors = [
            'h1',
            'meta[property="og:title"]',
            'meta[name="twitter:title"]',
            'title'
        ]
        
        for selector in title_selectors:
            try:
                element = page.locator(selector).first
                if element.count() > 0:
                    if selector == 'h1':
                        metadata['title'] = element.inner_text(timeout=2000).strip()
                    else:
                        metadata['title'] = element.get_attribute('content').strip()
                    if metadata['title']:
                        break
            except:
                continue
        
        # Author
        author_selectors = [
            'meta[name="author"]',
            'meta[property="article:author"]',
            '[rel="author"]',
            '.author',
            '[class*="author"]'
        ]
        
        for selector in author_selectors:
            try:
                element = page.locator(selector).first
                if element.count() > 0:
                    if 'meta' in selector:
                        metadata['author'] = element.get_attribute('content').strip()
                    else:
                        metadata['author'] = element.inner_text(timeout=2000).strip()
                    if metadata['author']:
                        break
            except:
                continue
        
        # Publish date
        date_selectors = [
            'meta[property="article:published_time"]',
            'meta[name="publish-date"]',
            'time[datetime]',
            '[class*="date"]'
        ]
        
        for selector in date_selectors:
            try:
                element = page.locator(selector).first
                if element.count() > 0:
                    if 'meta' in selector:
                        metadata['publish_date'] = element.get_attribute('content').strip()
                    else:
                        metadata['publish_date'] = element.get_attribute('datetime') or element.inner_text(timeout=2000)
                    if metadata['publish_date']:
                        break
            except:
                continue
        
        return metadata
    
    def _extract_domain(self, url: str) -> str:
        """Extract domain name from URL."""
        from urllib.parse import urlparse
        domain = urlparse(url).netloc
        return domain.replace('www.', '')

