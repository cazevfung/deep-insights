"""
Rich source ingestion service for URLs, file uploads, and manual text.
"""
from __future__ import annotations

import io
import json
import os
import shutil
import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from fastapi import UploadFile
from loguru import logger

try:
    import pandas as pd
except Exception:  # pragma: no cover - optional at import time
    pd = None

from utils.link_formatter import build_items, current_batch_id, iso_timestamp


# File size restrictions removed - no limit on upload count or size
MAX_TOTAL_BYTES = None  # No limit
MAX_SINGLE_FILE_BYTES = None  # No limit

TEXT_EXTENSIONS = {'.txt', '.md', '.markdown'}
WORD_EXTENSIONS = {'.doc', '.docx'}
PPT_EXTENSIONS = {'.ppt', '.pptx'}
EXCEL_EXTENSIONS = {'.xls', '.xlsx'}
PDF_EXTENSIONS = {'.pdf'}
SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | WORD_EXTENSIONS | PPT_EXTENSIONS | EXCEL_EXTENSIONS | PDF_EXTENSIONS


@dataclass
class IngestionItem:
    """Represents a normalized ingestion source mapped to a link entry."""

    url: str
    source_kind: str  # url | upload | text
    display_name: str
    metadata: Dict[str, Any]


class IngestionService:
    """Service to normalize mixed sources into scraper-compatible artifacts."""

    def __init__(self):
        self.project_root = Path(__file__).parent.parent.parent.parent
        self.links_file = self.project_root / "tests" / "data" / "test_links.json"
        self.upload_root = self.project_root / "backend" / "uploads"
        self.storage_root = self.project_root / "data" / "research"
        self.upload_root.mkdir(parents=True, exist_ok=True)

    async def ingest_sources(
        self,
        urls: List[str],
        text_entries: List[Dict[str, str]],
        files: List[UploadFile],
        session_id: Optional[str],
    ) -> Dict[str, Any]:
        """Normalize sources, update batch artifacts, and return batch metadata."""
        if not (urls or text_entries or files):
            raise ValueError("At least one link, upload, or text entry is required")

        logger.info(
            "Ingesting sources: %s urls, %s text entries, %s files",
            len(urls),
            len(text_entries),
            len(files),
        )

        # File size validation removed - no restrictions on upload size
        total_bytes = 0
        for upload in files:
            size = await self._get_upload_size(upload)
            total_bytes += size
        # No total size limit check

        batch_id = current_batch_id()
        target_session_id = session_id or batch_id

        batch_storage = self._ensure_batch_storage(target_session_id, batch_id)
        manifest_items: List[Dict[str, Any]] = []
        normalized_inputs: List[IngestionItem] = []

        # 1. URLs (unchanged)
        for raw_url in urls:
            url = raw_url.strip()
            if not url:
                continue
            normalized_inputs.append(
                IngestionItem(
                    url=url,
                    source_kind="url",
                    display_name=url,
                    metadata={"original_url": url},
                )
            )
            manifest_items.append(
                {
                    "source_kind": "url",
                    "display_name": url,
                    "original_url": url,
                }
            )

        # 2. Text entries
        for idx, entry in enumerate(text_entries, start=1):
            content = (entry.get("content") or "").strip()
            if not content:
                continue
            label = entry.get("label") or f"Text #{idx}"
            text_path = self._write_text_entry(content, label, batch_storage)
            normalized_inputs.append(
                IngestionItem(
                    url=text_path.as_uri(),
                    source_kind="text",
                    display_name=label,
                    metadata={
                        "label": label,
                        "character_count": len(content),
                        "file_path": str(text_path),
                    },
                )
            )
            manifest_items.append(
                {
                    "source_kind": "text",
                    "display_name": label,
                    "character_count": len(content),
                    "text_file": str(text_path),
                }
            )

        # 3. File uploads
        logger.info("Processing %s file upload(s)...", len(files))
        for idx, upload in enumerate(files, start=1):
            logger.info("Processing file %s/%s: %s", idx, len(files), upload.filename)
            try:
                normalized = await self._convert_upload(upload, batch_storage)
                normalized_inputs.append(normalized["item"])
                manifest_items.append(normalized["manifest"])
                logger.info("Successfully processed file %s/%s: %s", idx, len(files), upload.filename)
            except Exception as e:
                logger.error("Failed to process file %s/%s (%s): %s", idx, len(files), upload.filename, e, exc_info=True)
                raise ValueError(f"Failed to process file {upload.filename}: {str(e)}") from e

        if not normalized_inputs:
            raise ValueError("No valid sources were provided after filtering")

        # Build link entries compatible with scrapers
        ordered_urls = [entry.url for entry in normalized_inputs]
        link_items = build_items(ordered_urls)

        # Map manifest metadata with generated link IDs
        for item, manifest in zip(link_items, manifest_items):
            manifest["link_id"] = item["id"]
            manifest["link_type"] = item["type"]
            manifest["url"] = item["url"]

        payload = {
            "batchId": batch_id,
            "createdAt": iso_timestamp(),
            "links": link_items,
        }
        self._write_links_file(payload)
        self._write_manifest(batch_storage, target_session_id, batch_id, manifest_items)

        logger.info(
            "Prepared batch %s: %s items (%s urls, %s uploads, %s texts)",
            batch_id,
            len(link_items),
            len(urls),
            len(files),
            len(text_entries),
        )

        return {
            "batch_id": batch_id,
            "items": [
                {
                    "url": item["url"],
                    "source": item["type"],
                    "link_id": item["id"],
                }
                for item in link_items
            ],
            "total": len(link_items),
            "session_id": target_session_id,
            "manifest": manifest_items,
        }

    async def _convert_upload(self, upload: UploadFile, batch_storage: Path) -> Dict[str, Any]:
        """Convert uploaded file into normalized artifact and manifest.
        
        For PDFs and other files that need processing, we save the raw file
        and create a file:// URL. The file will be processed by the scraper
        workflow (ArticleScraper) alongside other sources, using the same
        worker pool.
        """
        filename = upload.filename or f"upload-{uuid.uuid4().hex}"
        suffix = Path(filename).suffix.lower()
        if suffix not in SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file type: {filename}")

        raw_dir = batch_storage / "uploads"
        raw_dir.mkdir(parents=True, exist_ok=True)
        raw_path = raw_dir / f"{uuid.uuid4().hex}{suffix}"

        # Persist raw bytes
        logger.info("Reading upload %s (%s bytes)...", filename, upload.size if hasattr(upload, 'size') else 'unknown')
        contents = await upload.read()
        raw_path.write_bytes(contents)
        logger.info("Saved upload %s to %s (%s bytes)", filename, raw_path, len(contents))

        # For text files (txt, md), extract immediately for backward compatibility
        # For PDFs, Word docs, Excel, etc. - save raw file and let scraper handle it
        if suffix in TEXT_EXTENSIONS:
            # Text files: extract immediately (fast, no processing needed)
            logger.info("Extracting text from text file %s...", filename)
            try:
                text_content = raw_path.read_text(encoding="utf-8", errors="ignore")
                if not text_content.strip():
                    raise ValueError(f"No textual content could be extracted from {filename}")
                
                # Save extracted text
                text_file = batch_storage / "prepared" / f"{raw_path.stem}.txt"
                text_file.parent.mkdir(parents=True, exist_ok=True)
                text_file.write_text(text_content, encoding="utf-8")
                
                item = IngestionItem(
                    url=text_file.as_uri(),
                    source_kind="upload",
                    display_name=filename,
                    metadata={
                        "original_name": filename,
                        "file_extension": suffix,
                        "raw_path": str(raw_path),
                        "text_path": str(text_file),
                    },
                )
                
                manifest = {
                    "source_kind": "upload",
                    "display_name": filename,
                    "raw_path": str(raw_path),
                    "text_file": str(text_file),
                    "file_extension": suffix,
                }
            except Exception as e:
                logger.error("Error extracting text from %s: %s", filename, e, exc_info=True)
                raise ValueError(f"Failed to extract text from {filename}: {str(e)}")
        else:
            # PDFs, Word, Excel, etc. - save raw file, let scraper process it
            logger.info("Saving %s for scraper workflow processing (type: %s)", filename, suffix)
            
            # Create file:// URL pointing to the raw file
            # Use absolute path for file:// URL
            file_url = raw_path.resolve().as_uri()
            
            item = IngestionItem(
                url=file_url,  # file:// URL - will be processed by ArticleScraper
                source_kind="upload",
                display_name=filename,
                metadata={
                    "original_name": filename,
                    "file_extension": suffix,
                    "raw_path": str(raw_path),
                    "processed_by": "scraper_workflow",  # Indicates this will be processed by scraper
                },
            )
            
            manifest = {
                "source_kind": "upload",
                "display_name": filename,
                "raw_path": str(raw_path),
                "file_extension": suffix,
                "processed_by": "scraper_workflow",
            }
            
            if suffix in EXCEL_EXTENSIONS and pd is not None:
                # For Excel, we can still extract a preview during ingestion
                try:
                    manifest["structured_preview"] = self._extract_excel_preview(raw_path)
                except Exception as e:
                    logger.warning("Failed to extract Excel preview: %s", e)

        return {"item": item, "manifest": manifest}

    def _extract_text_from_file(self, path: Path, suffix: str) -> str:
        """Route to the appropriate converter."""
        logger.debug("Extracting text from file: %s (suffix: %s)", path, suffix)
        if suffix in TEXT_EXTENSIONS:
            logger.debug("Using plain text reader for %s", path)
            return path.read_text(encoding="utf-8", errors="ignore")
        if suffix in PDF_EXTENSIONS:
            # Convert PDF to markdown format
            logger.info("Converting PDF to markdown: %s", path)
            try:
                from utils.pdf_to_markdown import convert_pdf_to_markdown
                logger.debug("Using PDF to markdown converter")
                result = convert_pdf_to_markdown(path)
                logger.info("PDF conversion completed: %s characters extracted", len(result))
                return result
            except ImportError as e:
                logger.warning(f"PDF to markdown converter not available: {e}. Falling back to plain text extraction.")
                # Fallback to plain text extraction
                from pypdf import PdfReader
                logger.debug("Using pypdf fallback")
                reader = PdfReader(str(path))
                pages = [page.extract_text() or "" for page in reader.pages]
                result = "\n\n".join(pages)
                logger.info("PDF text extraction completed: %s characters extracted", len(result))
                return result
            except Exception as e:
                logger.error(f"Error converting PDF to markdown: {e}. Falling back to plain text extraction.", exc_info=True)
                # Fallback to plain text extraction
                try:
                    from pypdf import PdfReader
                    logger.debug("Using pypdf fallback after error")
                    reader = PdfReader(str(path))
                    pages = [page.extract_text() or "" for page in reader.pages]
                    result = "\n\n".join(pages)
                    logger.info("PDF text extraction (fallback) completed: %s characters extracted", len(result))
                    return result
                except Exception as fallback_error:
                    logger.error(f"Fallback PDF extraction also failed: {fallback_error}", exc_info=True)
                    raise ValueError(f"Failed to extract text from PDF: {str(e)}") from e
        if suffix in WORD_EXTENSIONS:
            from docx import Document

            doc = Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs)
        if suffix in PPT_EXTENSIONS:
            from pptx import Presentation

            prs = Presentation(str(path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
            return "\n".join(texts)
        if suffix in EXCEL_EXTENSIONS:
            if pd is None:
                raise RuntimeError("pandas is required for Excel ingestion")
            data = self._extract_excel_preview(path)
            return json.dumps(data, ensure_ascii=False, indent=2)
        raise ValueError(f"Unsupported file extension: {suffix}")

    def _extract_excel_preview(self, path: Path) -> Dict[str, Any]:
        """Convert Excel sheets into JSON-ready structure."""
        if pd is None:
            raise RuntimeError("pandas is required for Excel ingestion")
        sheets = pd.read_excel(path, sheet_name=None, engine="openpyxl")
        preview = {}
        for sheet_name, df in sheets.items():
            sanitized = df.dropna(how="all")
            sanitized = sanitized.dropna(axis=1, how="all")
            sanitized = sanitized.fillna("")
            preview[str(sheet_name)] = {
                "columns": [str(col) for col in sanitized.columns],
                "records": sanitized.astype(str).to_dict(orient="records"),
            }
        return preview

    def _write_text_entry(self, content: str, label: str, batch_storage: Path) -> Path:
        """Persist inline text as a file for scraping."""
        safe_label = "".join(ch for ch in label if ch.isalnum() or ch in (" ", "_", "-")).strip().replace(" ", "_")
        if not safe_label:
            safe_label = uuid.uuid4().hex[:8]
        text_dir = batch_storage / "text"
        text_dir.mkdir(parents=True, exist_ok=True)
        path = text_dir / f"{safe_label}.txt"
        path.write_text(content, encoding="utf-8")
        return path

    def _write_links_file(self, payload: Dict[str, Any]) -> None:
        """Persist formatted links for downstream scrapers."""
        self.links_file.parent.mkdir(parents=True, exist_ok=True)
        self.links_file.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
        logger.info("Wrote %s links to %s", len(payload.get("links", [])), self.links_file)

    def _write_manifest(
        self,
        batch_storage: Path,
        session_id: str,
        batch_id: str,
        manifest_items: List[Dict[str, Any]],
    ) -> None:
        """Persist manifest alongside research artifacts."""
        manifest = {
            "session_id": session_id,
            "batch_id": batch_id,
            "created_at": iso_timestamp(),
            "items": manifest_items,
        }
        manifest_path = batch_storage / "manifest.json"
        manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
        logger.info("Manifest written to %s", manifest_path)

    def _ensure_batch_storage(self, session_id: str, batch_id: str) -> Path:
        """Ensure storage directory exists for session/batch artifacts."""
        base = self.storage_root / session_id / batch_id
        base.mkdir(parents=True, exist_ok=True)
        return base

    async def _get_upload_size(self, upload: UploadFile) -> int:
        """Determine upload size without losing the stream."""
        current = await upload.read()
        size = len(current)
        if hasattr(upload, "seek"):
            await upload.seek(0)
        else:  # pragma: no cover - fallback
            upload.file.seek(0)
        return size

